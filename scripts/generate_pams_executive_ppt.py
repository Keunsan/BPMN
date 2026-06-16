# -*- coding: utf-8 -*-
"""PAMS 임원 보고 PT 생성 — QnC 와이드 템플릿 기반"""

from __future__ import annotations

import shutil
from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# ── 경로 ──────────────────────────────────────────────
ROOT = Path(__file__).resolve().parents[1]
TEMPLATE = ROOT / "docs" / "PAMS_Executive_Report.pptx"
OUTPUT = ROOT / "docs" / "PAMS_Executive_Report_Final.pptx"

# ── 브랜드 컬러 (템플릿 theme 기준) ───────────────────
C_PRIMARY = RGBColor(0x00, 0x70, 0xC0)
C_PRIMARY_DK = RGBColor(0x05, 0x63, 0xC1)
C_ACCENT = RGBColor(0xED, 0x7D, 0x31)
C_GREEN = RGBColor(0x70, 0xAD, 0x47)
C_DARK = RGBColor(0x44, 0x54, 0x6A)
C_GRAY = RGBColor(0xA5, 0xA5, 0xA5)
C_LIGHT = RGBColor(0xF2, 0xF6, 0xFA)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_RED = RGBColor(0xC0, 0x00, 0x00)
C_GOLD = RGBColor(0xFF, 0xC0, 0x00)

FONT = "맑은 고딕"
TODAY = date.today().strftime("%Y.%m.%d")


# ── 유틸 ──────────────────────────────────────────────
def emu_in(v: float):
    return Inches(v)


def set_shape_fill(shape, color: RGBColor, transparency: float = 0.0):
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color
    if transparency:
        fill.fore_color.brightness = transparency


def add_shadow(shape, blur: int = 40000, dist: int = 20000):
    sp = shape._element.spPr
    if sp.effectLst is None:
        from pptx.oxml.ns import qn

        effectLst = sp.makeelement(qn("a:effectLst"), {})
        sp.append(effectLst)
    effectLst = sp.effectLst
    outer = effectLst.find("{http://schemas.openxmlformats.org/drawingml/2006/main}outerShdw")
    if outer is None:
        from pptx.oxml.ns import qn

        outer = effectLst.makeelement(
            qn("a:outerShdw"),
            {
                "blurRad": str(blur),
                "dist": str(dist),
                "dir": "5400000",
                "algn": "bl",
                "rotWithShape": "0",
            },
        )
        srgb = outer.makeelement(
            "{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr",
            {"val": "000000"},
        )
        alpha = srgb.makeelement(
            "{http://schemas.openxmlformats.org/drawingml/2006/main}alpha",
            {"val": "25000"},
        )
        srgb.append(alpha)
        outer.append(srgb)
        effectLst.append(outer)


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    *,
    size: int = 14,
    bold: bool = False,
    color: RGBColor = C_DARK,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    line_spacing: float = 1.25,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.space_after = Pt(4)
    p.line_spacing = line_spacing
    run = p.runs[0] if p.runs else p.add_run()
    if not p.runs:
        run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(
    slide,
    left,
    top,
    width,
    height,
    items: list[str],
    *,
    size: int = 13,
    color: RGBColor = C_DARK,
    bullet_color: RGBColor = C_PRIMARY,
    spacing: int = 6,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(spacing)
        p.line_spacing = 1.2
        if p.runs:
            run = p.runs[0]
        else:
            run = p.add_run()
            run.text = item
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.color.rgb = color
        pPr = p._pPr
        if pPr is None:
            from pptx.oxml.ns import qn

            pPr = p._p.get_or_add_pPr()
        # bullet char
        from pptx.oxml.ns import qn

        buChar = pPr.makeelement(qn("a:buChar"), {"char": "•"})
        buClr = pPr.makeelement(qn("a:buClr"))
        hex_val = str(bullet_color)
        srgb = buClr.makeelement(qn("a:srgbClr"), {"val": hex_val})
        buClr.append(srgb)
        for old in pPr.findall(qn("a:buChar")):
            pPr.remove(old)
        for old in pPr.findall(qn("a:buClr")):
            pPr.remove(old)
        pPr.insert(0, buClr)
        pPr.insert(1, buChar)
    return box


def add_rounded_card(
    slide,
    left,
    top,
    width,
    height,
    fill: RGBColor = C_WHITE,
    line: RGBColor | None = RGBColor(0xE0, 0xE8, 0xF0),
):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    set_shape_fill(shape, fill)
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    add_shadow(shape)
    return shape


def add_accent_bar(slide, left, top, height, color: RGBColor = C_PRIMARY, width=emu_in(0.06)):
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height
    )
    set_shape_fill(bar, color)
    bar.line.fill.background()
    return bar


def add_circle_badge(slide, left, top, size, text: str, fill: RGBColor = C_PRIMARY):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL, left, top, size, size
    )
    set_shape_fill(shape, fill)
    shape.line.fill.background()
    tf = shape.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = C_WHITE
    return shape


def set_title(slide, title: str, subtitle: str | None = None):
    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:
            ph.text = title
            for p in ph.text_frame.paragraphs:
                for run in p.runs:
                    run.font.name = FONT
        elif idx == 1 and subtitle:
            ph.text = subtitle
            for p in ph.text_frame.paragraphs:
                for run in p.runs:
                    run.font.name = FONT
                    run.font.size = Pt(16)
                    run.font.color.rgb = C_GRAY


def set_footer_date(slide):
    for ph in slide.placeholders:
        if ph.placeholder_format.type == 16:  # DATE
            ph.text = TODAY


def add_key_message(slide, text: str, top=emu_in(1.05)):
    card = add_rounded_card(
        slide, emu_in(0.45), top, emu_in(12.1), emu_in(0.65), fill=C_PRIMARY
    )
    card.line.fill.background()
    add_textbox(
        slide,
        emu_in(0.65),
        top + emu_in(0.12),
        emu_in(11.7),
        emu_in(0.45),
        text,
        size=15,
        bold=True,
        color=C_WHITE,
        align=PP_ALIGN.LEFT,
    )


def add_slide_header(slide, section_num: str, title: str):
    set_title(slide, title)
    set_footer_date(slide)
    badge = add_circle_badge(
        slide, emu_in(11.6), emu_in(0.15), emu_in(0.45), section_num, fill=C_ACCENT
    )
    return badge


def add_comparison_table(slide, left, top, headers, rows, col_widths):
    rows_n = len(rows) + 1
    cols_n = len(headers)
    table_shape = slide.shapes.add_table(
        rows_n, cols_n, left, top, sum(col_widths), emu_in(0.38 * rows_n)
    )
    table = table_shape.table
    for ci, w in enumerate(col_widths):
        table.columns[ci].width = w
    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_PRIMARY_DK
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.name = FONT
                run.font.size = Pt(11)
                run.font.bold = True
                run.font.color.rgb = C_WHITE
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val)
            if ri % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_LIGHT
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.name = FONT
                    run.font.size = Pt(10)
                    run.font.color.rgb = C_DARK
    return table_shape


def add_progress_bar(slide, left, top, width, pct: float, label: str):
    h = emu_in(0.22)
    bg = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, h
    )
    set_shape_fill(bg, RGBColor(0xE8, 0xEE, 0xF4))
    bg.line.fill.background()
    fill_w = int(width * min(max(pct, 0), 1))
    if fill_w > 0:
        fg = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, fill_w, h
        )
        set_shape_fill(fg, C_PRIMARY)
        fg.line.fill.background()
    add_textbox(
        slide, left, top - emu_in(0.28), width, emu_in(0.25),
        f"{label}  {int(pct * 100)}%",
        size=11, bold=True, color=C_DARK,
    )


def add_flow_arrow(slide, x1, y, x2, color=C_PRIMARY):
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, x1, y, x2, y
    )
    conn.line.color.rgb = color
    conn.line.width = Pt(2)
    # arrow head via triangle
    tri = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE,
        x2 - emu_in(0.12), y - emu_in(0.08), emu_in(0.14), emu_in(0.16),
    )
    set_shape_fill(tri, color)
    tri.line.fill.background()
    tri.rotation = 90.0


# ── 슬라이드 빌더 ─────────────────────────────────────
def build_cover(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_title(
        slide,
        "PAMS 구축 현황 및 향후 추진 방향",
        "Process Architecture Management System\n프로세스 아키텍처 관리 시스템",
    )
    # decorative elements
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, emu_in(4.8), prs.slide_width, emu_in(0.08)
    )
    set_shape_fill(bar, C_PRIMARY)
    bar.line.fill.background()

    deco = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, emu_in(0.45), emu_in(3.6), emu_in(2.5), emu_in(0.06)
    )
    set_shape_fill(deco, C_ACCENT)
    deco.line.fill.background()

    add_textbox(
        slide, emu_in(0.45), emu_in(5.05), emu_in(4), emu_in(0.35),
        f"보고일  {TODAY}", size=12, color=C_GRAY,
    )
    add_textbox(
        slide, emu_in(0.45), emu_in(5.4), emu_in(6), emu_in(0.35),
        "정보전략 / 프로세스 혁신", size=13, bold=True, color=C_DARK,
    )

    # decorative translucent circles
    for x, y, r, col in [
        (10.5, 0.8, 1.2, C_PRIMARY),
        (11.2, 2.0, 0.7, C_ACCENT),
        (9.8, 1.5, 0.5, C_GREEN),
    ]:
        c = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL, emu_in(x), emu_in(y), emu_in(r), emu_in(r)
        )
        set_shape_fill(c, col)
        c.line.fill.background()
        c.fill.transparency = 0.88

    set_footer_date(slide)


def build_toc(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[11])
    set_title(slide, "목  차")
    set_footer_date(slide)
    add_accent_bar(slide, emu_in(0.45), emu_in(1.15), emu_in(5.5))

    items = [
        ("01", "보고 개요 — 왜 PAMS인가"),
        ("02", "해결하고자 하는 과제"),
        ("03", "PAMS 비전과 차별점"),
        ("04", "구축 현황 — Phase별 진행"),
        ("05", "핵심 구현 성과"),
        ("06", "현재 위치 — 완료 vs 잔여"),
        ("07", "향후 전략 방향"),
        ("08", "추진 로드맵"),
        ("09", "기대 효과"),
        ("10", "결론 및 제안"),
    ]
    for i, (num, label) in enumerate(items):
        col = i // 5
        row = i % 5
        x = emu_in(0.65 + col * 6.2)
        y = emu_in(1.35 + row * 1.05)
        add_rounded_card(slide, x, y, emu_in(5.8), emu_in(0.82), fill=C_LIGHT if row % 2 else C_WHITE)
        add_circle_badge(slide, x + emu_in(0.15), y + emu_in(0.18), emu_in(0.48), num)
        add_textbox(
            slide, x + emu_in(0.75), y + emu_in(0.22), emu_in(4.8), emu_in(0.4),
            label, size=13, bold=True, color=C_DARK,
        )


def build_slide01(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[11])
    add_slide_header(slide, "01", "보고 개요 — 왜 PAMS인가")
    add_key_message(
        slide,
        "BPMN 도구가 아닌, 업무·시스템·데이터·거버넌스를 하나로 연결하는 전사 운영 플랫폼",
    )
    cards = [
        ("E2E 표준화", "전사 공통 표준과\n법인·사업부별 변형 관리", C_PRIMARY),
        ("변경 영향 분석", "시스템·데이터·업무\n연쇄 영향 즉시 파악", C_ACCENT),
        ("운영 거버넌스", "승인·버전·이력\n감사 대응 체계화", C_GREEN),
    ]
    for i, (title, desc, col) in enumerate(cards):
        x = emu_in(0.45 + i * 4.15)
        card = add_rounded_card(slide, x, emu_in(2.0), emu_in(3.85), emu_in(2.2), fill=C_WHITE)
        add_accent_bar(slide, x, emu_in(2.0), emu_in(2.2), color=col)
        add_textbox(slide, x + emu_in(0.25), emu_in(2.15), emu_in(3.3), emu_in(0.4),
                    title, size=16, bold=True, color=col)
        add_textbox(slide, x + emu_in(0.25), emu_in(2.65), emu_in(3.3), emu_in(1.2),
                    desc, size=12, color=C_DARK)
    add_bullets(
        slide, emu_in(0.55), emu_in(4.5), emu_in(12), emu_in(1.5),
        [
            "기존 업무 흐름도 도구는 '그리기·보관'에 그치며, ERP/MES 등 운영 시스템과의 연계가 어렵습니다.",
            "PAMS는 한국어·영어·중국어(번체)를 지원하여 해외 법인과의 협업 기반을 마련합니다.",
        ],
        size=12, bullet_color=C_PRIMARY,
    )


def build_slide02(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[11])
    add_slide_header(slide, "02", "해결하고자 하는 과제")
    add_key_message(slide, "프로세스·시스템·데이터·조직 정보가 분산되어 있어, 변경·장애·감사 대응에 시간과 리스크가 큽니다.")

    pains = [
        ("조직", "사업부·공장·법인별 업무 방식 상이\n'표준' 정의 부재", C_PRIMARY),
        ("시스템", "ERP/MES 변경 시 연관 업무 추적 불가\n장애 시 영향 범위 파악 지연", C_ACCENT),
        ("거버넌스", "문서 최신성 미보장\n변경 이력·승인 체계·감사 자료 분산", C_RED),
        ("글로벌", "해외 법인과의\n언어·프로세스 커뮤니케이션 장벽", C_GREEN),
    ]
    for i, (title, desc, col) in enumerate(pains):
        x = emu_in(0.45 + (i % 2) * 6.3)
        y = emu_in(1.95 + (i // 2) * 1.85)
        card = add_rounded_card(slide, x, y, emu_in(5.95), emu_in(1.55), fill=C_WHITE)
        add_circle_badge(slide, x + emu_in(0.2), y + emu_in(0.35), emu_in(0.55), title[:1], fill=col)
        add_textbox(slide, x + emu_in(0.9), y + emu_in(0.2), emu_in(4.8), emu_in(0.35),
                    title, size=14, bold=True, color=col)
        add_textbox(slide, x + emu_in(0.9), y + emu_in(0.6), emu_in(4.8), emu_in(0.8),
                    desc, size=11, color=C_DARK)

    # KPI targets
    kpi_card = add_rounded_card(slide, emu_in(0.45), emu_in(5.65), emu_in(12.1), emu_in(1.15), fill=C_PRIMARY_DK)
    kpi_card.line.fill.background()
    kpis = ["영향 분석 70%↓", "감사 준비 80%↓", "장애 대응 50%↓", "표준 준수 90%+"]
    for i, k in enumerate(kpis):
        add_textbox(
            slide, emu_in(0.7 + i * 3.0), emu_in(5.95), emu_in(2.8), emu_in(0.5),
            k, size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER,
        )


def build_slide03(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[11])
    add_slide_header(slide, "03", "PAMS 비전과 차별점")
    add_key_message(
        slide,
        "전사 업무가 어떻게 흐르고, 어떤 시스템·데이터를 쓰는지 한곳에서 보고 변경 영향을 즉시 판단",
    )

    headers = ["구분", "기존 업무 흐름도 도구", "PAMS"]
    rows = [
        ["목적", "다이어그램 작성·저장", "전사 운영관리 플랫폼"],
        ["범위", "개별 프로세스", "L1~L4 + E2E 가치흐름"],
        ["연계", "독립 운영", "ERP/MES/SCM 메타데이터 통합"],
        ["활용", "문서 아카이브", "영향도·변경관리·감사"],
    ]
    add_comparison_table(
        slide, emu_in(0.45), emu_in(1.85), headers, rows,
        [emu_in(1.5), emu_in(5.0), emu_in(5.6)],
    )

    # 3-layer diagram
    layers = [
        ("Layer A", "프로세스 계층", "L1~L4 · E2E 카탈로그", C_PRIMARY),
        ("Layer B", "업무 흐름", "L3 상세 · E2E 전사 흐름", C_ACCENT),
        ("Layer C", "운영 정보", "Task · 시스템 · 데이터 · RACI", C_GREEN),
    ]
    for i, (tag, name, desc, col) in enumerate(layers):
        y = emu_in(4.05 + i * 0.72)
        box = add_rounded_card(slide, emu_in(0.55), y, emu_in(12), emu_in(0.58), fill=C_LIGHT)
        add_textbox(slide, emu_in(0.75), y + emu_in(0.08), emu_in(1.2), emu_in(0.4),
                    tag, size=10, bold=True, color=col)
        add_textbox(slide, emu_in(2.0), y + emu_in(0.08), emu_in(2.5), emu_in(0.4),
                    name, size=12, bold=True, color=C_DARK)
        add_textbox(slide, emu_in(4.5), y + emu_in(0.08), emu_in(7.5), emu_in(0.4),
                    desc, size=11, color=C_GRAY)
        if i < 2:
            ay = y + emu_in(0.58)
            conn = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT, emu_in(6.5), ay, emu_in(6.5), ay + emu_in(0.14)
            )
            conn.line.color.rgb = col
            conn.line.width = Pt(2)


def build_slide04(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[11])
    add_slide_header(slide, "04", "구축 현황 — Phase별 진행")
    add_key_message(slide, "2026년 6월 기준, 핵심 기반 기능 70% 이상 구현 완료")

    phases = [
        ("0~1", "기술 기반·화면 골격", True),
        ("2", "L1~L4 프로세스·승인·버전", True),
        ("2+", "전사 표준 / 현장 변형", True),
        ("3", "업무 흐름도·버전비교", True),
        ("4", "Task·시스템·데이터 연결", True),
        ("5", "운영 지식그래프", True),
        ("7", "E2E 프로세스 카탈로그", True),
        ("4~6", "RACI·대시보드·배포", False),
    ]
    for i, (phase, label, done) in enumerate(phases):
        col = i // 4
        row = i % 4
        x = emu_in(0.45 + col * 6.3)
        y = emu_in(1.85 + row * 0.95)
        fill = C_GREEN if done else RGBColor(0xFF, 0xF0, 0xE0)
        border = C_GREEN if done else C_ACCENT
        card = add_rounded_card(slide, x, y, emu_in(5.95), emu_in(0.78), fill=fill)
        card.line.color.rgb = border
        icon = "✓" if done else "…"
        add_circle_badge(slide, x + emu_in(0.12), y + emu_in(0.12), emu_in(0.52), icon,
                         fill=C_GREEN if done else C_ACCENT)
        add_textbox(slide, x + emu_in(0.75), y + emu_in(0.1), emu_in(0.8), emu_in(0.35),
                    f"P{phase}", size=10, bold=True, color=C_DARK)
        add_textbox(slide, x + emu_in(1.5), y + emu_in(0.15), emu_in(4.2), emu_in(0.5),
                    label, size=11, color=C_DARK)

    add_progress_bar(slide, emu_in(0.55), emu_in(5.85), emu_in(11.9), 0.72, "전체 구축 진행률")
    add_textbox(
        slide, emu_in(0.55), emu_in(6.35), emu_in(12), emu_in(0.4),
        "DB 설계 25건 · 화면·기능 40여 개 · 3개국어 지원",
        size=11, color=C_GRAY,
    )


def build_slide05(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[11])
    add_slide_header(slide, "05", "핵심 구현 성과")
    add_key_message(slide, "단순 등록·조회를 넘어, 실제 운영 시나리오를 염두에 둔 차별 기능이 동작합니다.")

    achievements = [
        ("①", "표준 vs 변형", "법인·사업부별 변형 관리\n표준과 변형 비교", C_PRIMARY),
        ("②", "시스템·데이터 연동", "ERP/MES 테이블 실시간 조회\nTask↔시스템↔데이터 3단 연결", C_ACCENT),
        ("③", "E2E 카탈로그", "전사 가치흐름(Source to Pay 등)\nE2E→L3→L4 상세 펼치기", C_GREEN),
        ("④", "운영 지식그래프", "프로세스·시스템·데이터\n관계를 그래프로 탐색", C_PRIMARY_DK),
        ("⑤", "자동 동기화", "업무 흐름도 선행 관계\n→ Task 선행 업무 자동 반영", C_GOLD),
    ]
    for i, (num, title, desc, col) in enumerate(achievements):
        col_i = i % 3
        row_i = i // 3
        x = emu_in(0.45 + col_i * 4.15)
        y = emu_in(1.85 + row_i * 2.15)
        w = emu_in(3.85)
        h = emu_in(1.85)
        card = add_rounded_card(slide, x, y, w, h)
        add_accent_bar(slide, x, y, h, color=col)
        add_textbox(slide, x + emu_in(0.2), y + emu_in(0.15), emu_in(0.5), emu_in(0.35),
                    num, size=18, bold=True, color=col)
        add_textbox(slide, x + emu_in(0.2), y + emu_in(0.55), emu_in(3.4), emu_in(0.4),
                    title, size=13, bold=True, color=C_DARK)
        add_textbox(slide, x + emu_in(0.2), y + emu_in(1.0), emu_in(3.4), emu_in(0.7),
                    desc, size=10, color=C_GRAY)


def build_slide06(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[11])
    add_slide_header(slide, "06", "현재 위치 — 완료 vs 잔여")
    add_key_message(slide, "'데이터를 모으는 플랫폼'은 거의 완성 → '모아둔 데이터로 의사결정' 단계로 진입")

    done_items = [
        "프로세스·업무 흐름도·E2E",
        "Task·시스템·데이터 연결",
        "운영 지식그래프 (탐색)",
        "승인 대기함 (기본)",
        "외부 시스템·화면 마스터",
        "표준/변형 비교",
    ]
    todo_items = [
        "RACI / KPI·리스크·통제",
        "조직·역할·사용자 관리",
        "영향도·통합검색·히트맵",
        "변경 이력·감사 로그",
        "대시보드 실데이터",
        "로그인·권한 본격 연동",
    ]

    for label, items, col, x in [
        ("완료 ✅", done_items, C_GREEN, emu_in(0.45)),
        ("잔여 📋", todo_items, C_ACCENT, emu_in(6.55)),
    ]:
        card = add_rounded_card(slide, x, emu_in(1.85), emu_in(5.95), emu_in(3.8), fill=C_WHITE)
        hdr = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, emu_in(1.85), emu_in(5.95), emu_in(0.45)
        )
        set_shape_fill(hdr, col)
        hdr.line.fill.background()
        add_textbox(slide, x + emu_in(0.2), emu_in(1.9), emu_in(5.5), emu_in(0.35),
                    label, size=14, bold=True, color=C_WHITE)
        for j, item in enumerate(items):
            dot = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.OVAL,
                x + emu_in(0.25), emu_in(2.5 + j * 0.48), emu_in(0.12), emu_in(0.12),
            )
            set_shape_fill(dot, col)
            dot.line.fill.background()
            add_textbox(
                slide, x + emu_in(0.5), emu_in(2.42 + j * 0.48), emu_in(5.2), emu_in(0.4),
                item, size=11, color=C_DARK,
            )

    add_textbox(
        slide, emu_in(0.55), emu_in(5.9), emu_in(12), emu_in(0.5),
        "현재 Phase 4~5 중반  |  RACI·KPI·조직 등 DB 설계는 완료, 화면·기능 연결 잔여",
        size=11, bold=True, color=C_PRIMARY,
    )


def build_slide07(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[11])
    add_slide_header(slide, "07", "향후 전략 방향")
    add_key_message(slide, "PRD 잔여 과제 완성 + ITSM·로그 등 운영 데이터 연계로 PAMS 가치 확대")

    themes = [
        ("A", "운영 지식그래프\n완성", "영향·사용·변경\n실데이터 연결", C_PRIMARY),
        ("B", "거버넌스\n규칙", "승인 전 필수 검증\n변경 이력·감사", C_ACCENT),
        ("C", "표준/변형\n거버넌스", "메타 비교\n준수율 가시화", C_GREEN),
        ("D", "분석·Layer C\n완성", "RACI/KPI·대시보드\n영향도·검색", C_PRIMARY_DK),
        ("E", "운영 데이터\n연계 ★", "ITSM VOC\n시스템 활용도", C_RED),
    ]
    for i, (tag, title, desc, col) in enumerate(themes):
        x = emu_in(0.35 + i * 2.55)
        card = add_rounded_card(slide, x, emu_in(1.85), emu_in(2.35), emu_in(2.5))
        add_circle_badge(slide, x + emu_in(0.85), emu_in(2.0), emu_in(0.55), tag, fill=col)
        add_textbox(slide, x + emu_in(0.15), emu_in(2.65), emu_in(2.05), emu_in(0.7),
                    title, size=11, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
        add_textbox(slide, x + emu_in(0.15), emu_in(3.35), emu_in(2.05), emu_in(0.8),
                    desc, size=9, color=C_GRAY, align=PP_ALIGN.CENTER)

    # E theme detail box
    detail = add_rounded_card(slide, emu_in(0.45), emu_in(4.65), emu_in(12.1), emu_in(2.2), fill=C_LIGHT)
    add_textbox(slide, emu_in(0.65), emu_in(4.75), emu_in(11.5), emu_in(0.35),
                "운영 데이터 연계 (2차 확장)", size=13, bold=True, color=C_RED)
    add_bullets(
        slide, emu_in(0.65), emu_in(5.15), emu_in(5.5), emu_in(1.5),
        [
            "ITSM VOC: PAMS 프로세스 트리를 ITSM에서 참조·매핑",
            "PAMS가 매핑 결과를 수집 → 프로세스별 SR 현황 분석",
            "E2E·법인별 VOC 추이로 개선 우선순위 도출",
        ],
        size=10, bullet_color=C_RED,
    )
    add_bullets(
        slide, emu_in(6.5), emu_in(5.15), emu_in(5.5), emu_in(1.5),
        [
            "메뉴별 권한·접속 로그 → 화면·Task·L3까지 집계",
            "정의된 화면 vs 실사용 Gap 분석",
            "VOC(多) × 활용도(低) = 최우선 개선 후보",
        ],
        size=10, bullet_color=C_PRIMARY,
    )


def build_slide08(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[11])
    add_slide_header(slide, "08", "추진 로드맵")
    add_key_message(slide, "가치 실현 → Layer C → 거버넌스 → 운영 연계 → 확장 순으로 추진")

    steps = [
        ("1순위", "가치 실현", "영향도·대시보드\n변경 이력", C_PRIMARY, True),
        ("2순위", "Layer C", "RACI/KPI\n필수항목 검증", C_ACCENT, True),
        ("3순위", "거버넌스", "로그인·권한\nLevel별 승인", C_GREEN, False),
        ("4순위", "운영 연계 ★", "ITSM VOC\n활용도 Heat Map", C_RED, False),
        ("5순위", "확장", "AI 보조\n전사 배포", C_GRAY, False),
    ]
    for i, (rank, title, desc, col, active) in enumerate(steps):
        x = emu_in(0.35 + i * 2.55)
        w = emu_in(2.35)
        card = add_rounded_card(
            slide, x, emu_in(1.85), w, emu_in(2.3),
            fill=col if active else C_WHITE,
        )
        if not active:
            card.line.color.rgb = col
        tc = C_WHITE if active else col
        add_textbox(slide, x + emu_in(0.1), emu_in(2.0), w, emu_in(0.35),
                    rank, size=10, bold=True, color=tc, align=PP_ALIGN.CENTER)
        add_textbox(slide, x + emu_in(0.1), emu_in(2.4), w, emu_in(0.45),
                    title, size=12, bold=True, color=tc, align=PP_ALIGN.CENTER)
        add_textbox(slide, x + emu_in(0.1), emu_in(3.0), w, emu_in(0.9),
                    desc, size=9, color=tc if active else C_GRAY, align=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            add_flow_arrow(slide, x + w, emu_in(3.0), x + w + emu_in(0.18), col)

    # ITSM flow
    flow_y = emu_in(4.55)
    boxes = [
        ("PAMS\n프로세스 트리", C_PRIMARY),
        ("ITSM\nVOC 매핑", C_ACCENT),
        ("PAMS\n분석·보고", C_GREEN),
    ]
    for i, (txt, col) in enumerate(boxes):
        x = emu_in(1.5 + i * 4.0)
        box = add_rounded_card(slide, x, flow_y, emu_in(2.8), emu_in(0.9), fill=col)
        box.line.fill.background()
        add_textbox(slide, x, flow_y + emu_in(0.15), emu_in(2.8), emu_in(0.65),
                    txt, size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        if i < 2:
            add_flow_arrow(slide, x + emu_in(2.85), flow_y + emu_in(0.45),
                          x + emu_in(3.95), C_DARK)

    add_textbox(
        slide, emu_in(0.55), emu_in(5.75), emu_in(12), emu_in(0.8),
        "파일럿: Source to Pay(구매-입고) 1개 영역  |  ITSM VOC 매핑 + ERP/MM 메뉴 활용도",
        size=11, bold=True, color=C_DARK, align=PP_ALIGN.CENTER,
    )


def build_slide09(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[11])
    add_slide_header(slide, "09", "기대 효과와 성공 지표")
    add_key_message(slide, "운영 효율·리스크 관리·감사 대응·글로벌 협업에서 정량적 효과 기대")

    metrics = [
        ("표준화", "90%+", "전사 표준 + 변형 관리", C_PRIMARY),
        ("변경 관리", "70%↓", "영향 분석 시간 단축", C_ACCENT),
        ("장애 대응", "50%↓", "연관 업무 즉시 추적", C_GREEN),
        ("감사·통제", "80%↓", "자료 준비 시간 단축", C_PRIMARY_DK),
    ]
    for i, (area, target, desc, col) in enumerate(metrics):
        x = emu_in(0.45 + i * 3.15)
        card = add_rounded_card(slide, x, emu_in(1.85), emu_in(2.9), emu_in(1.8))
        ring = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.DONUT, x + emu_in(0.85), emu_in(2.0), emu_in(1.2), emu_in(1.2)
        )
        set_shape_fill(ring, col)
        ring.line.fill.background()
        inner = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            x + emu_in(1.05), emu_in(2.2), emu_in(0.8), emu_in(0.8),
        )
        set_shape_fill(inner, C_WHITE)
        inner.line.fill.background()
        add_textbox(slide, x + emu_in(0.85), emu_in(2.35), emu_in(1.2), emu_in(0.5),
                    target, size=14, bold=True, color=col, align=PP_ALIGN.CENTER)
        add_textbox(slide, x + emu_in(0.15), emu_in(3.35), emu_in(2.6), emu_in(0.35),
                    area, size=12, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
        add_textbox(slide, x + emu_in(0.15), emu_in(3.65), emu_in(2.6), emu_in(0.35),
                    desc, size=9, color=C_GRAY, align=PP_ALIGN.CENTER)

    add_rounded_card(slide, emu_in(0.45), emu_in(4.0), emu_in(12.1), emu_in(1.5), fill=C_LIGHT)
    add_textbox(slide, emu_in(0.65), emu_in(4.1), emu_in(11.5), emu_in(0.35),
                "2차 확장(운영 연계) 추가 효과", size=12, bold=True, color=C_RED)
    add_bullets(
        slide, emu_in(0.65), emu_in(4.5), emu_in(11.5), emu_in(0.9),
        [
            "ITSM VOC → L3/E2E 단위 SR 현황으로 데이터 기반 개선·교육 우선순위",
            "활용도 분석 → 미사용 화면·Ghost process 식별, 통합·권한 정리",
            "SR(多) × 활용도(低) 영역 = 최우선 개선 후보 자동 도출",
        ],
        size=10, bullet_color=C_RED,
    )


def build_slide10(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[11])
    add_slide_header(slide, "10", "결론 및 제안")
    add_key_message(slide, "PAMS 1차(정의 플랫폼) 핵심 기반 완료 → 2차(분석·운영 연계)로 현업 ROI 증명")

    summary_card = add_rounded_card(slide, emu_in(0.45), emu_in(1.85), emu_in(12.1), emu_in(1.35), fill=C_PRIMARY_DK)
    summary_card.line.fill.background()
    add_bullets(
        slide, emu_in(0.65), emu_in(1.95), emu_in(11.5), emu_in(1.2),
        [
            "✅  L1~L4 + E2E + 업무 흐름도 + Task·시스템·데이터 + 운영 지식그래프 기반 완성",
            "✅  표준/변형, 외부 시스템 연동, E2E 상세 펼치기 등 차별 기능 구현",
            "📋  RACI/KPI·대시보드·영향도 + ITSM·활용도 연계 잔여",
        ],
        size=11, color=C_WHITE, bullet_color=C_GOLD,
    )

    proposals = [
        ("1", "Source to Pay 파일럿", "프로세스 + ITSM VOC + MM 활용도"),
        ("2", "1순위 2~3개월 집중", "영향도·대시보드·변경이력 → 체감 ROI"),
        ("3", "Process Owner·ITSM 협의", "VOC L3/L4 매핑 · PAMS = 전사 기준"),
        ("4", "권한 연동 후 전사 확대", "로그인·역할 기반 접근 제어 선행"),
    ]
    for i, (num, title, desc) in enumerate(proposals):
        x = emu_in(0.45 + (i % 2) * 6.3)
        y = emu_in(3.5 + (i // 2) * 1.35)
        card = add_rounded_card(slide, x, y, emu_in(5.95), emu_in(1.1))
        add_circle_badge(slide, x + emu_in(0.15), y + emu_in(0.25), emu_in(0.55), num, fill=C_ACCENT)
        add_textbox(slide, x + emu_in(0.85), y + emu_in(0.2), emu_in(4.8), emu_in(0.35),
                    title, size=12, bold=True, color=C_DARK)
        add_textbox(slide, x + emu_in(0.85), y + emu_in(0.6), emu_in(4.8), emu_in(0.4),
                    desc, size=10, color=C_GRAY)

    closing = add_rounded_card(slide, emu_in(0.45), emu_in(6.2), emu_in(12.1), emu_in(0.65), fill=C_ACCENT)
    closing.line.fill.background()
    add_textbox(
        slide, emu_in(0.65), emu_in(6.32), emu_in(11.7), emu_in(0.45),
        "PAMS는 프로세스를 '문서'가 아니라 '운영 자산'으로 관리하는 전사 디지털 전환의 핵심 인프라입니다.",
        size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER,
    )


def remove_all_slides(prs: Presentation):
    """템플릿 기본 슬라이드를 관계까지 제거한다."""
    while len(prs.slides) > 0:
        r_id = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(r_id)
        del prs.slides._sldIdLst[0]


def main():
    if not TEMPLATE.exists():
        raise FileNotFoundError(f"Template not found: {TEMPLATE}")

    shutil.copy2(TEMPLATE, OUTPUT)
    prs = Presentation(str(OUTPUT))
    remove_all_slides(prs)

    build_cover(prs)
    build_toc(prs)
    build_slide01(prs)
    build_slide02(prs)
    build_slide03(prs)
    build_slide04(prs)
    build_slide05(prs)
    build_slide06(prs)
    build_slide07(prs)
    build_slide08(prs)
    build_slide09(prs)
    build_slide10(prs)

    prs.save(str(OUTPUT))
    print(f"Generated: {OUTPUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
